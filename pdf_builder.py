"""
pdf_builder.py
Consumes content dict from content_generator.py.
Produces 3 merged PDFs:
  - Standard Pupil  (Voc + Ret + Inf, 7 questions each)
  - Supported Pupil (Voc + Ret + Inf, 5 questions each, with scaffolds)
  - All Answers     (6 pages: Std + Sup for each lesson)

Renders all 9 question formats:
  open_line, find_and_copy, numbered_list,
  tick_one, tick_two, true_false_table,
  sequencing, reason_evidence_table, two_part_ab
"""

import os
import hashlib
import random
from io import BytesIO

from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.lib.units import mm
from pypdf import PdfReader, PdfWriter

from wfa_shared.brand import year_colour_rgb, hex_to_rgb, WFA_BLUE, YEAR_COLOURS

# ---------------------------------------------------------------------------
# Page geometry
# ---------------------------------------------------------------------------
W, H = A4
MARGIN = 8 * mm
CW = W - 2 * MARGIN        # usable content width
MIN_Y = 14 * mm             # bottom margin — nothing drawn below this

# ---------------------------------------------------------------------------
# Colours
# ---------------------------------------------------------------------------
BOX_BORDER = (0.173, 0.173, 0.424)   # #2c2c6c — borders, question labels
BOX_BG     = (0.941, 0.941, 0.973)   # #f0f0f8 — reading text box
GREEN      = (0.102, 0.478, 0.102)   # #1a7a1a — answer text and highlights
DARK       = (0.133, 0.133, 0.133)   # body text
GREY_LINE  = (0.6,   0.6,   0.6  )  # ruled lines, cell borders
LIGHT_GREY = (0.94,  0.94,  0.94 )  # table cell backgrounds
TICK_CELL  = (0.85,  0.95,  0.85 )  # answer highlight cell

# Year group brand colours — sourced from wfa-shared
# Active colour — updated by set_year_group(); defaults to Y4
_ACCENT = year_colour_rgb("Y4")


def set_year_group(year_group: str):
    """Update the active accent colour for PDF branding. Call before build_pdfs()."""
    global _ACCENT
    _ACCENT = year_colour_rgb(year_group)


def _coerce_answer(q):
    """
    Return the answer field as a plain string.
    Guards against Claude occasionally returning a list instead of a string.
    """
    val = q.get('answer', '')
    if isinstance(val, list):
        return ' / '.join(str(v) for v in val)
    return str(val) if val is not None else ''


# ===========================================================================
# Low-level drawing utilities
# ===========================================================================

def wrap_text(c, text, font, size, max_w):
    """Return list of lines that fit within max_w."""
    words = str(text).split()
    lines, line = [], ''
    for w in words:
        test = (line + ' ' + w).strip()
        if c.stringWidth(test, font, size) <= max_w:
            line = test
        else:
            if line:
                lines.append(line)
            line = w
    if line:
        lines.append(line)
    return lines or ['']


def text_height(c, text, font, size, max_w, line_leading=1.35):
    """Return height of wrapped text block in points."""
    lines = wrap_text(c, text, font, size, max_w)
    return len(lines) * size * line_leading


def draw_wrapped(c, text, font, size, x, y, max_w, line_leading=1.35):
    """Draw wrapped text starting at (x, y). Returns y after last line."""
    c.setFont(font, size)
    lines = wrap_text(c, text, font, size, max_w)
    for line in lines:
        c.drawString(x, y, line)
        y -= size * line_leading
    return y


def draw_checkbox(c, x, y, size=4 * mm, filled=False, fill_colour=None):
    """Draw a small square checkbox. Returns (right_edge, centre_y)."""
    if filled and fill_colour:
        c.setFillColorRGB(*fill_colour)
        c.rect(x, y - size, size, size, fill=1, stroke=0)
    c.setFillColorRGB(1, 1, 1) if not filled else None
    c.setStrokeColorRGB(*GREY_LINE)
    c.setLineWidth(0.4)
    c.rect(x, y - size, size, size, fill=0, stroke=1)
    return x + size, y - size / 2


def marks_label(c, marks, y):
    """Print 'N mark(s)' right-aligned at current y."""
    label = f"{marks} mark" if marks == 1 else f"{marks} marks"
    c.setFont("Helvetica-Oblique", 7)
    c.setFillColorRGB(0.5, 0.5, 0.5)
    c.drawRightString(MARGIN + CW, y, label)


def draw_tick_instruction(c, instruction, y):
    """Print 'Tick one.' or 'Tick two.' right-aligned."""
    c.setFont("Helvetica-BoldOblique", 8)
    c.setFillColorRGB(*BOX_BORDER)
    c.drawRightString(MARGIN + CW, y, instruction)
    return y - 4 * mm


def draw_ruled_line(c, y, gap=6.5 * mm):
    """Draw one answer line at y. Returns y below line."""
    c.setStrokeColorRGB(*GREY_LINE)
    c.setLineWidth(0.4)
    c.line(MARGIN, y - gap, MARGIN + CW, y - gap)
    return y - gap


def deterministic_shuffle(items):
    """Shuffle a list reproducibly based on its own content."""
    seed = int(hashlib.md5('|'.join(items).encode()).hexdigest()[:8], 16)
    r = random.Random(seed)
    result = list(items)
    r.shuffle(result)
    return result


# ===========================================================================
# Header
# ===========================================================================

def draw_header(c, lesson_type, day, date, key_q, i_can_statements, icon_path):
    """Draw the learning label. Returns y immediately below the final rule."""
    y = H - MARGIN

    # Row 1: "Learning Focus"  [icon]  "Day Date"
    c.setFont("Helvetica-Bold", 8)
    c.setFillColorRGB(*DARK)
    c.drawString(MARGIN, y - 5 * mm, "Learning Focus")

    icon_x = MARGIN + 32 * mm
    try:
        c.drawImage(icon_path, icon_x, y - 7 * mm,
                    width=7 * mm, height=7 * mm,
                    mask='auto', preserveAspectRatio=True)
    except Exception:
        pass

    c.setFont("Helvetica", 8)
    c.drawString(icon_x + 8 * mm, y - 5 * mm, f"{day}  {date}")
    y -= 8 * mm

    # Learning focus text — bold, underlined, year group colour
    c.setFont("Helvetica-Bold", 10)
    c.setFillColorRGB(*_ACCENT)
    c.drawString(MARGIN, y - 4 * mm, key_q)
    kq_w = c.stringWidth(key_q, "Helvetica-Bold", 10)
    c.setLineWidth(0.5)
    c.setStrokeColorRGB(*_ACCENT)
    c.line(MARGIN, y - 5 * mm, MARGIN + kq_w, y - 5 * mm)
    y -= 7 * mm

    # I can statements
    c.setFont("Helvetica", 8)
    c.setFillColorRGB(*DARK)
    for ican in i_can_statements:
        c.drawString(MARGIN, y - 3.5 * mm, ican)
        y -= 4.5 * mm

    # Rule
    y -= 1 * mm
    c.setStrokeColorRGB(*GREY_LINE)
    c.setLineWidth(0.5)
    c.line(MARGIN, y, MARGIN + CW, y)
    y -= 2.5 * mm

    return y


# ===========================================================================
# Reading text box
# ===========================================================================

def draw_text_box(c, text, y_top, font_size=10):
    """Draw the reading extract box. Returns y below box."""
    lines = wrap_text(c, text, "Helvetica", font_size, CW - 6 * mm)
    lh = font_size * 1.4
    box_h = len(lines) * lh + 6 * mm

    c.setFillColorRGB(*BOX_BG)
    c.setStrokeColorRGB(*BOX_BORDER)
    c.setLineWidth(0.8)
    c.roundRect(MARGIN, y_top - box_h, CW, box_h, 2 * mm, fill=1, stroke=1)

    c.setFillColorRGB(*DARK)
    c.setFont("Helvetica", font_size)
    ty = y_top - 3.5 * mm - font_size * 0.72
    for line in lines:
        c.drawString(MARGIN + 3 * mm, ty, line)
        ty -= lh

    return y_top - box_h - 3 * mm


# ===========================================================================
# Question preamble (text_reference + question label)
# ===========================================================================

def draw_preamble(c, q, y):
    ref = q.get('text_reference', '').strip()
    if ref:
        c.setFont("Helvetica-Oblique", 7.5)
        c.setFillColorRGB(0.5, 0.5, 0.5)
        c.drawString(MARGIN, y, ref)
        y -= 7.5 * 1.3

    c.setFillColorRGB(*DARK)
    c.setFont("Helvetica-Bold", 9)
    label = f"{q['number']}. "
    lw = c.stringWidth(label, "Helvetica-Bold", 9)
    c.drawString(MARGIN, y, label)

    qtext = q.get('question', '')
    parts = qtext.split('\n\n', 1)
    if len(parts) == 2:
        quote, question_body = parts
        c.setFont("Helvetica-BoldOblique", 9)
        c.drawString(MARGIN + lw, y, quote)
        y -= 9 * 1.35
        c.setFont("Helvetica-Bold", 9)
        q_lines = wrap_text(c, question_body, "Helvetica-Bold", 9, CW - lw)
        for i, line in enumerate(q_lines):
            c.drawString(MARGIN + (lw if i == 0 else 0), y, line)
            y -= 9 * 1.35
    else:
        q_lines = wrap_text(c, qtext, "Helvetica-Bold", 9, CW - lw)
        for i, line in enumerate(q_lines):
            c.drawString(MARGIN + (lw if i == 0 else 0), y, line)
            y -= 9 * 1.35

    y -= 1 * mm
    return y, lw


def draw_scaffold(c, scaffold, y):
    if not scaffold:
        return y
    c.setFont("Helvetica-Oblique", 8)
    c.setFillColorRGB(0.4, 0.4, 0.6)
    c.drawString(MARGIN + 3 * mm, y, scaffold)
    y -= 8 * 1.35
    return y


# ===========================================================================
# Format renderers — pupil versions
# ===========================================================================

def render_open_line_pupil(c, q, y, is_supported):
    lines_n = q['format_data'].get('lines', 2)
    if is_supported:
        y = draw_scaffold(c, q.get('supported_scaffold'), y)
    gap = 6.5 * mm
    c.setStrokeColorRGB(*GREY_LINE)
    c.setLineWidth(0.4)
    for i in range(lines_n):
        ly = y - (i + 1) * gap
        c.line(MARGIN, ly, MARGIN + CW, ly)
    return y - lines_n * gap - 2 * mm


def render_find_and_copy_pupil(c, q, y, is_supported):
    c.setFont("Helvetica", 8)
    c.setFillColorRGB(*DARK)
    c.drawString(MARGIN, y - 3 * mm, "_______________________________________________")
    return y - 8 * mm


def render_numbered_list_pupil(c, q, y, is_supported):
    n = q['format_data'].get('num_points', 2)
    if is_supported:
        y = draw_scaffold(c, q.get('supported_scaffold'), y)
    gap = 6.5 * mm
    c.setStrokeColorRGB(*GREY_LINE)
    c.setLineWidth(0.4)
    c.setFont("Helvetica-Bold", 9)
    c.setFillColorRGB(*DARK)
    for i in range(n):
        c.drawString(MARGIN, y - (i * (gap * 2)) - 3 * mm, f"{i + 1}.")
        c.line(MARGIN + 5 * mm, y - (i * (gap * 2)) - gap, MARGIN + CW, y - (i * (gap * 2)) - gap)
        c.line(MARGIN + 5 * mm, y - (i * (gap * 2)) - gap * 2, MARGIN + CW, y - (i * (gap * 2)) - gap * 2)
    return y - n * gap * 2 - 2 * mm


def render_tick_one_pupil(c, q, y):
    fd = q['format_data']
    options = fd.get('options', [])
    y = draw_tick_instruction(c, "Tick one.", y)
    row_h = 6.5 * mm
    col_w = CW / 2
    c.setStrokeColorRGB(*GREY_LINE)
    c.setLineWidth(0.4)
    for row in range(2):
        for col in range(2):
            idx = row * 2 + col
            if idx >= len(options):
                break
            x = MARGIN + col * col_w
            ry = y - row * row_h
            c.setFillColorRGB(1, 1, 1)
            c.rect(x, ry - row_h, col_w, row_h, fill=1, stroke=1)
            c.setFillColorRGB(*DARK)
            c.setFont("Helvetica", 8.5)
            c.drawString(x + 2 * mm, ry - row_h + 2 * mm, options[idx])
    return y - 2 * row_h - 2 * mm


def render_tick_two_pupil(c, q, y):
    fd = q['format_data']
    options = fd.get('options', [])
    y = draw_tick_instruction(c, "Tick two.", y)
    row_h = 6 * mm
    c.setStrokeColorRGB(*GREY_LINE)
    c.setLineWidth(0.4)
    for i, opt in enumerate(options):
        ry = y - i * row_h
        c.setFillColorRGB(1, 1, 1)
        c.rect(MARGIN, ry - row_h, CW, row_h, fill=1, stroke=1)
        draw_checkbox(c, MARGIN + 2 * mm, ry - row_h * 0.2)
        c.setFillColorRGB(*DARK)
        c.setFont("Helvetica", 8.5)
        c.drawString(MARGIN + 8 * mm, ry - row_h + 2 * mm, opt)
    return y - len(options) * row_h - 2 * mm


def render_true_false_table_pupil(c, q, y):
    fd = q['format_data']
    statements = fd.get('statements', [])
    col_stmt = CW * 0.72
    col_tf = CW * 0.14
    row_h = 7 * mm

    c.setFillColorRGB(*BOX_BORDER)
    c.rect(MARGIN, y - row_h, CW, row_h, fill=1, stroke=0)
    c.setFillColorRGB(1, 1, 1)
    c.setFont("Helvetica-Bold", 8.5)
    c.drawString(MARGIN + 2 * mm, y - row_h + 2 * mm, "Statement")
    c.drawCentredString(MARGIN + col_stmt + col_tf / 2, y - row_h + 2 * mm, "True")
    c.drawCentredString(MARGIN + col_stmt + col_tf + col_tf / 2, y - row_h + 2 * mm, "False")
    y -= row_h

    c.setStrokeColorRGB(*GREY_LINE)
    c.setLineWidth(0.4)
    for i, stmt in enumerate(statements):
        fill = LIGHT_GREY if i % 2 == 0 else (1, 1, 1)
        c.setFillColorRGB(*fill)
        c.rect(MARGIN, y - row_h, CW, row_h, fill=1, stroke=1)
        c.setFillColorRGB(*DARK)
        c.setFont("Helvetica", 8)
        c.drawString(MARGIN + 2 * mm, y - row_h + 2 * mm, stmt.get('text', ''))
        c.setLineWidth(0.3)
        c.line(MARGIN + col_stmt, y, MARGIN + col_stmt, y - row_h)
        c.line(MARGIN + col_stmt + col_tf, y, MARGIN + col_stmt + col_tf, y - row_h)
        y -= row_h

    return y - 2 * mm


def render_sequencing_pupil(c, q, y):
    fd = q['format_data']
    items = deterministic_shuffle(fd.get('items', []))
    row_h = 6 * mm
    box_sz = 4.5 * mm

    c.setFont("Helvetica", 8.5)
    c.setFillColorRGB(*DARK)
    for i, item in enumerate(items):
        ry = y - i * row_h
        c.setFillColorRGB(1, 1, 1)
        c.setStrokeColorRGB(*BOX_BORDER)
        c.setLineWidth(0.5)
        c.rect(MARGIN, ry - box_sz, box_sz, box_sz, fill=1, stroke=1)
        c.setFillColorRGB(*DARK)
        c.setFont("Helvetica", 8.5)
        c.drawString(MARGIN + box_sz + 2 * mm, ry - row_h + 2 * mm, item)
    return y - len(items) * row_h - 2 * mm


def render_reason_evidence_pupil(c, q, y):
    fd = q['format_data']
    example = fd.get('example', {})
    blank_rows = fd.get('rows', 2)
    col_r = CW * 0.42
    col_e = CW - col_r
    header_h = 6 * mm
    row_h = 11 * mm

    c.setFillColorRGB(*BOX_BORDER)
    c.rect(MARGIN, y - header_h, col_r, header_h, fill=1, stroke=0)
    c.rect(MARGIN + col_r, y - header_h, col_e, header_h, fill=1, stroke=0)
    c.setFillColorRGB(1, 1, 1)
    c.setFont("Helvetica-Bold", 8.5)
    c.drawCentredString(MARGIN + col_r / 2, y - header_h + 2 * mm, "Reason")
    c.drawCentredString(MARGIN + col_r + col_e / 2, y - header_h + 2 * mm, "Evidence")
    y -= header_h

    c.setFillColorRGB(*LIGHT_GREY)
    c.setStrokeColorRGB(*GREY_LINE)
    c.setLineWidth(0.4)
    c.rect(MARGIN, y - row_h, col_r, row_h, fill=1, stroke=1)
    c.rect(MARGIN + col_r, y - row_h, col_e, row_h, fill=1, stroke=1)
    c.setFillColorRGB(0.35, 0.35, 0.35)
    c.setFont("Helvetica-Oblique", 8)
    reason_lines = wrap_text(c, example.get('reason', ''), "Helvetica-Oblique", 8, col_r - 4 * mm)
    ry = y - 3 * mm
    for line in reason_lines:
        c.drawString(MARGIN + 2 * mm, ry, line)
        ry -= 8 * 1.3
    ev_lines = wrap_text(c, example.get('evidence', ''), "Helvetica-Oblique", 8, col_e - 4 * mm)
    ry = y - 3 * mm
    for line in ev_lines:
        c.drawString(MARGIN + col_r + 2 * mm, ry, line)
        ry -= 8 * 1.3
    y -= row_h

    for _ in range(blank_rows):
        c.setFillColorRGB(1, 1, 1)
        c.setStrokeColorRGB(*GREY_LINE)
        c.setLineWidth(0.4)
        c.rect(MARGIN, y - row_h, col_r, row_h, fill=1, stroke=1)
        c.rect(MARGIN + col_r, y - row_h, col_e, row_h, fill=1, stroke=1)
        y -= row_h

    return y - 2 * mm


def render_two_part_ab_pupil(c, q, y, is_supported):
    fd = q['format_data']
    parts = fd.get('parts', [])
    for part in parts:
        label = part.get('label', '?')
        pq = part.get('question', '')
        c.setFont("Helvetica-Bold", 8.5)
        c.setFillColorRGB(*DARK)
        pl = f"({label})  "
        plw = c.stringWidth(pl, "Helvetica-Bold", 8.5)
        c.drawString(MARGIN + 3 * mm, y, pl)
        pq_lines = wrap_text(c, pq, "Helvetica-Bold", 8.5, CW - plw - 3 * mm)
        for i, line in enumerate(pq_lines):
            c.drawString(MARGIN + 3 * mm + plw, y - i * (8.5 * 1.35), line)
        y -= len(pq_lines) * (8.5 * 1.35) + 1 * mm
        c.setStrokeColorRGB(*GREY_LINE)
        c.setLineWidth(0.4)
        c.line(MARGIN + 3 * mm, y - 6 * mm, MARGIN + CW, y - 6 * mm)
        c.line(MARGIN + 3 * mm, y - 12 * mm, MARGIN + CW, y - 12 * mm)
        y -= 13 * mm
    return y - 2 * mm


# ===========================================================================
# Format renderers — answer versions
# ===========================================================================

def render_open_line_answer(c, q, y):
    answer = _coerce_answer(q)
    c.setFillColorRGB(*GREEN)
    c.setFont("Helvetica-Oblique", 8.5)
    ans_lines = wrap_text(c, answer, "Helvetica-Oblique", 8.5, CW)
    for line in ans_lines:
        c.drawString(MARGIN, y - 5 * mm, line)
        y -= 5 * mm
    return y - 3 * mm


def render_find_and_copy_answer(c, q, y):
    target = q['format_data'].get('target_word', _coerce_answer(q))
    c.setFont("Helvetica-BoldOblique", 9)
    c.setFillColorRGB(*GREEN)
    c.drawString(MARGIN, y - 4 * mm, f"\u2714  {target}")
    return y - 9 * mm


def render_numbered_list_answer(c, q, y):
    answer = _coerce_answer(q)
    n = q['format_data'].get('num_points', 2)
    import re
    matches = re.split(r'\d+\.\s+', answer)
    parts = [p.strip() for p in matches if p.strip()]
    if len(parts) < n:
        parts = [answer] + [''] * (n - 1)
    c.setFont("Helvetica-Oblique", 8.5)
    c.setFillColorRGB(*GREEN)
    for i in range(n):
        c.setFont("Helvetica-Bold", 9)
        c.setFillColorRGB(*DARK)
        c.drawString(MARGIN, y - 3 * mm, f"{i + 1}.")
        c.setFont("Helvetica-Oblique", 8.5)
        c.setFillColorRGB(*GREEN)
        c.drawString(MARGIN + 5 * mm, y - 3 * mm, parts[i] if i < len(parts) else '')
        y -= 9 * mm
    return y - 2 * mm


def render_tick_one_answer(c, q, y):
    fd = q['format_data']
    options = fd.get('options', [])
    correct_i = fd.get('correct_index', 0)
    y = draw_tick_instruction(c, "Tick one.", y)
    row_h = 6.5 * mm
    col_w = CW / 2
    c.setStrokeColorRGB(*GREY_LINE)
    c.setLineWidth(0.4)
    for row in range(2):
        for col in range(2):
            idx = row * 2 + col
            if idx >= len(options):
                break
            x = MARGIN + col * col_w
            ry = y - row * row_h
            is_correct = (idx == correct_i)
            c.setFillColorRGB(*TICK_CELL) if is_correct else c.setFillColorRGB(1, 1, 1)
            c.rect(x, ry - row_h, col_w, row_h, fill=1, stroke=1)
            if is_correct:
                c.setFillColorRGB(*GREEN)
                c.setFont("Helvetica-Bold", 8.5)
                c.drawString(x + 2 * mm, ry - row_h + 2 * mm, options[idx] + "  \u2714")
            else:
                c.setFillColorRGB(*DARK)
                c.setFont("Helvetica", 8.5)
                c.drawString(x + 2 * mm, ry - row_h + 2 * mm, options[idx])
    return y - 2 * row_h - 2 * mm


def render_tick_two_answer(c, q, y):
    fd = q['format_data']
    options = fd.get('options', [])
    correct_is = set(fd.get('correct_indices', []))
    y = draw_tick_instruction(c, "Tick two.", y)
    row_h = 6 * mm
    c.setStrokeColorRGB(*GREY_LINE)
    c.setLineWidth(0.4)
    for i, opt in enumerate(options):
        ry = y - i * row_h
        is_correct = (i in correct_is)
        fill = TICK_CELL if is_correct else (1, 1, 1)
        c.setFillColorRGB(*fill)
        c.rect(MARGIN, ry - row_h, CW, row_h, fill=1, stroke=1)
        if is_correct:
            c.setFillColorRGB(*GREEN)
            c.setFont("Helvetica-Bold", 8.5)
            c.drawString(MARGIN + 8 * mm, ry - row_h + 2 * mm, opt + "  \u2714")
        else:
            c.setFillColorRGB(*DARK)
            c.setFont("Helvetica", 8.5)
            c.drawString(MARGIN + 8 * mm, ry - row_h + 2 * mm, opt)
    return y - len(options) * row_h - 2 * mm


def render_true_false_table_answer(c, q, y):
    fd = q['format_data']
    statements = fd.get('statements', [])
    col_stmt = CW * 0.72
    col_tf = CW * 0.14
    row_h = 7 * mm

    c.setFillColorRGB(*BOX_BORDER)
    c.rect(MARGIN, y - row_h, CW, row_h, fill=1, stroke=0)
    c.setFillColorRGB(1, 1, 1)
    c.setFont("Helvetica-Bold", 8.5)
    c.drawString(MARGIN + 2 * mm, y - row_h + 2 * mm, "Statement")
    c.drawCentredString(MARGIN + col_stmt + col_tf / 2, y - row_h + 2 * mm, "True")
    c.drawCentredString(MARGIN + col_stmt + col_tf + col_tf / 2, y - row_h + 2 * mm, "False")
    y -= row_h

    c.setStrokeColorRGB(*GREY_LINE)
    c.setLineWidth(0.4)
    for i, stmt in enumerate(statements):
        fill = LIGHT_GREY if i % 2 == 0 else (1, 1, 1)
        c.setFillColorRGB(*fill)
        c.rect(MARGIN, y - row_h, CW, row_h, fill=1, stroke=1)
        c.setFillColorRGB(*DARK)
        c.setFont("Helvetica", 8)
        c.drawString(MARGIN + 2 * mm, y - row_h + 2 * mm, stmt.get('text', ''))
        c.line(MARGIN + col_stmt, y, MARGIN + col_stmt, y - row_h)
        c.line(MARGIN + col_stmt + col_tf, y, MARGIN + col_stmt + col_tf, y - row_h)
        is_true = stmt.get('correct', False)
        tick_x = (MARGIN + col_stmt + col_tf / 2) if is_true else (MARGIN + col_stmt + col_tf + col_tf / 2)
        c.setFillColorRGB(*GREEN)
        c.setFont("Helvetica-Bold", 9)
        c.drawCentredString(tick_x, y - row_h + 2 * mm, "\u2714")
        y -= row_h

    return y - 2 * mm


def render_sequencing_answer(c, q, y):
    fd = q['format_data']
    correct_order = fd.get('items', [])
    shuffled = deterministic_shuffle(correct_order)
    position = {item: str(i + 1) for i, item in enumerate(correct_order)}
    row_h = 6 * mm
    box_sz = 4.5 * mm

    for i, item in enumerate(shuffled):
        ry = y - i * row_h
        c.setFillColorRGB(*TICK_CELL)
        c.setStrokeColorRGB(*BOX_BORDER)
        c.setLineWidth(0.5)
        c.rect(MARGIN, ry - box_sz, box_sz, box_sz, fill=1, stroke=1)
        c.setFillColorRGB(*GREEN)
        c.setFont("Helvetica-Bold", 8.5)
        c.drawCentredString(MARGIN + box_sz / 2, ry - box_sz + 1 * mm, position[item])
        c.setFillColorRGB(*DARK)
        c.setFont("Helvetica", 8.5)
        c.drawString(MARGIN + box_sz + 2 * mm, ry - row_h + 2 * mm, item)
    return y - len(shuffled) * row_h - 2 * mm


def render_reason_evidence_answer(c, q, y):
    import re as _re

    fd = q['format_data']
    example = fd.get('example', {})
    blank_rows = fd.get('rows', 2)
    answer = _coerce_answer(q)
    col_r = CW * 0.42
    col_e = CW - col_r
    header_h = 6 * mm
    row_h = 13 * mm

    ans_pairs = []
    for raw in answer.strip().split('\n'):
        raw = raw.strip()
        if '|' in raw:
            r, e = raw.split('|', 1)
            ans_pairs.append((r.strip(), e.strip()))
    while len(ans_pairs) < blank_rows:
        ans_pairs.append(('', ''))

    c.setFillColorRGB(*BOX_BORDER)
    c.rect(MARGIN, y - header_h, col_r, header_h, fill=1, stroke=0)
    c.rect(MARGIN + col_r, y - header_h, col_e, header_h, fill=1, stroke=0)
    c.setFillColorRGB(1, 1, 1)
    c.setFont("Helvetica-Bold", 8.5)
    c.drawCentredString(MARGIN + col_r / 2, y - header_h + 2 * mm, "Reason")
    c.drawCentredString(MARGIN + col_r + col_e / 2, y - header_h + 2 * mm, "Evidence")
    y -= header_h

    c.setFillColorRGB(*LIGHT_GREY)
    c.setStrokeColorRGB(*GREY_LINE)
    c.setLineWidth(0.4)
    c.rect(MARGIN, y - row_h, col_r, row_h, fill=1, stroke=1)
    c.rect(MARGIN + col_r, y - row_h, col_e, row_h, fill=1, stroke=1)
    c.setFillColorRGB(0.35, 0.35, 0.35)
    c.setFont("Helvetica-Oblique", 8)
    ry = y - 3 * mm
    for line in wrap_text(c, example.get('reason', ''), "Helvetica-Oblique", 8, col_r - 4 * mm):
        c.drawString(MARGIN + 2 * mm, ry, line)
        ry -= 8 * 1.3
    ry = y - 3 * mm
    for line in wrap_text(c, example.get('evidence', ''), "Helvetica-Oblique", 8, col_e - 4 * mm):
        c.drawString(MARGIN + col_r + 2 * mm, ry, line)
        ry -= 8 * 1.3
    y -= row_h

    for i in range(blank_rows):
        c.setFillColorRGB(*TICK_CELL)
        c.setStrokeColorRGB(*GREY_LINE)
        c.setLineWidth(0.4)
        c.rect(MARGIN, y - row_h, col_r, row_h, fill=1, stroke=1)
        c.rect(MARGIN + col_r, y - row_h, col_e, row_h, fill=1, stroke=1)

        reason_text = ans_pairs[i][0] if i < len(ans_pairs) else ''
        evid_text   = ans_pairs[i][1] if i < len(ans_pairs) else ''

        if reason_text or evid_text:
            c.setFillColorRGB(*GREEN)
            c.setFont("Helvetica-Oblique", 8)
            ry = y - 3 * mm
            for line in wrap_text(c, reason_text, "Helvetica-Oblique", 8, col_r - 4 * mm):
                c.drawString(MARGIN + 2 * mm, ry, line)
                ry -= 8 * 1.3
            ry = y - 3 * mm
            for line in wrap_text(c, evid_text, "Helvetica-Oblique", 8, col_e - 4 * mm):
                c.drawString(MARGIN + col_r + 2 * mm, ry, line)
                ry -= 8 * 1.3
        else:
            c.setFillColorRGB(0.7, 0.7, 0.7)
            c.setFont("Helvetica-Oblique", 7.5)
            c.drawString(MARGIN + 2 * mm, y - row_h + 3 * mm, "See mark scheme")

        y -= row_h

    return y - 2 * mm


def render_two_part_ab_answer(c, q, y):
    fd = q['format_data']
    parts = fd.get('parts', [])
    for part in parts:
        label = part.get('label', '?')
        pq = part.get('question', '')
        ans = part.get('answer', '')
        c.setFont("Helvetica-Bold", 8.5)
        c.setFillColorRGB(*DARK)
        pl = f"({label})  "
        plw = c.stringWidth(pl, "Helvetica-Bold", 8.5)
        c.drawString(MARGIN + 3 * mm, y, pl)
        pq_lines = wrap_text(c, pq, "Helvetica-Bold", 8.5, CW - plw - 3 * mm)
        for i, line in enumerate(pq_lines):
            c.drawString(MARGIN + 3 * mm + plw, y - i * (8.5 * 1.35), line)
        y -= len(pq_lines) * (8.5 * 1.35) + 1 * mm
        c.setFont("Helvetica-Oblique", 8.5)
        c.setFillColorRGB(*GREEN)
        c.drawString(MARGIN + 5 * mm, y - 3 * mm, f"\u2714  {ans}")
        y -= 9 * mm
    return y - 2 * mm


# ---------------------------------------------------------------------------
# Draw-lines-matching renderers (KS1 only)
# ---------------------------------------------------------------------------

_DLM_COL_W  = None   # computed per call
_DLM_GAP_W  = None
_DLM_BOX_H  = 9 * mm
_DLM_DOT_R  = 1.2 * mm


def _dlm_geometry():
    col_w = CW * 0.43
    gap_w = CW * 0.14
    return col_w, gap_w


def _dlm_draw_boxes(c, left_items, right_items, y_top, draw_right_filled=False):
    """Draw the two columns of boxes. Returns list of (left_dot_x, left_dot_y) and
    (right_dot_x, right_dot_y) per row, plus final y."""
    col_w, gap_w = _dlm_geometry()
    box_h = _DLM_BOX_H
    r_col_x = MARGIN + col_w + gap_w

    left_dots  = []
    right_dots = []
    y = y_top

    n_rows = max(len(left_items), len(right_items))
    for i in range(n_rows):
        # Left box
        if i < len(left_items):
            c.setStrokeColorRGB(*GREY_LINE)
            c.setLineWidth(0.4)
            c.setFillColorRGB(1, 1, 1)
            c.rect(MARGIN, y - box_h, col_w, box_h, fill=1, stroke=1)
            c.setFillColorRGB(*DARK)
            c.setFont("Helvetica", 8)
            txt = left_items[i]
            t_lines = wrap_text(c, txt, "Helvetica", 8, col_w - 4 * mm)
            line_h = 8 * 1.3
            t_top = y - (box_h - len(t_lines) * line_h) / 2 - line_h * 0.15
            for j, ln in enumerate(t_lines):
                c.drawString(MARGIN + 2 * mm, t_top - j * line_h, ln)

        # Left dot (right edge of left box, vertically centred)
        ld_x = MARGIN + col_w
        ld_y = y - box_h / 2
        c.setFillColorRGB(*DARK)
        c.circle(ld_x, ld_y, _DLM_DOT_R, fill=1, stroke=0)
        left_dots.append((ld_x, ld_y))

        # Right box
        if i < len(right_items):
            c.setStrokeColorRGB(*GREY_LINE)
            c.setLineWidth(0.4)
            fill_col = TICK_CELL if draw_right_filled else (1, 1, 1)
            c.setFillColorRGB(*fill_col)
            c.rect(r_col_x, y - box_h, col_w, box_h, fill=1, stroke=1)
            c.setFillColorRGB(*DARK if not draw_right_filled else GREEN)
            c.setFont("Helvetica", 8)
            txt = right_items[i]
            t_lines = wrap_text(c, txt, "Helvetica", 8, col_w - 4 * mm)
            t_top = y - (box_h - len(t_lines) * 8 * 1.3) / 2 - 8 * 1.3 * 0.15
            for j, ln in enumerate(t_lines):
                c.drawString(r_col_x + 2 * mm, t_top - j * 8 * 1.3, ln)

        # Right dot (left edge of right box, vertically centred)
        rd_x = r_col_x
        rd_y = y - box_h / 2
        c.setFillColorRGB(*DARK)
        c.circle(rd_x, rd_y, _DLM_DOT_R, fill=1, stroke=0)
        right_dots.append((rd_x, rd_y))

        y -= box_h + 1.5 * mm

    return left_dots, right_dots, y


def render_draw_lines_matching_pupil(c, q, y):
    """Pupil version — boxes and dots, no lines drawn."""
    fd = q['format_data']
    left_items  = fd.get('left_items', [])
    right_items = fd.get('right_items', [])
    _, _, y = _dlm_draw_boxes(c, left_items, right_items, y)
    return y - 2 * mm


def render_draw_lines_matching_answer(c, q, y):
    """Answer/mark scheme version — boxes, dots, and correct lines in green."""
    fd = q['format_data']
    left_items   = fd.get('left_items', [])
    right_items  = fd.get('right_items', [])
    correct_pairs = fd.get('correct_pairs', [])

    left_dots, right_dots, y_end = _dlm_draw_boxes(
        c, left_items, right_items, y, draw_right_filled=False
    )

    # Draw connecting lines
    c.setStrokeColorRGB(*GREEN)
    c.setLineWidth(1.2)
    for (li, ri) in correct_pairs:
        if li < len(left_dots) and ri < len(right_dots):
            lx, ly = left_dots[li]
            rx, ry = right_dots[ri]
            c.line(lx, ly, rx, ry)

    # Tick mark beside each right box
    col_w, gap_w = _dlm_geometry()
    r_col_x = MARGIN + col_w + gap_w
    for (li, ri) in correct_pairs:
        if ri < len(right_dots):
            _, ry = right_dots[ri]
            c.setFont("Helvetica-Bold", 8)
            c.setFillColorRGB(*GREEN)
            c.drawString(MARGIN + CW - 6 * mm, ry - 3, "\u2714")

    return y_end - 2 * mm


# ===========================================================================
# Height estimation (for pupil/answer overflow checking)
# ===========================================================================

def estimate_height(c, q):
    ref_h = 7.5 * 1.3 if q.get('text_reference', '').strip() else 0
    qtext = q.get('question', '')
    lw = c.stringWidth(f"{q['number']}. ", "Helvetica-Bold", 9)
    q_lines = wrap_text(c, qtext.split('\n\n')[-1], "Helvetica-Bold", 9, CW - lw)
    q_text_h = len(q_lines) * 9 * 1.35 + 1 * mm

    fmt = q.get('format', 'open_line')
    fd = q.get('format_data', {})

    if fmt == 'open_line':
        body_h = fd.get('lines', 2) * 6.5 * mm + 2 * mm
    elif fmt == 'find_and_copy':
        body_h = 8 * mm
    elif fmt == 'numbered_list':
        body_h = fd.get('num_points', 2) * 13 * mm + 2 * mm
    elif fmt in ('tick_one',):
        body_h = 4 * mm + 2 * 6.5 * mm + 2 * mm
    elif fmt == 'tick_two':
        body_h = 4 * mm + len(fd.get('options', [])) * 6 * mm + 2 * mm
    elif fmt == 'true_false_table':
        body_h = (len(fd.get('statements', [])) + 1) * 7 * mm + 2 * mm
    elif fmt == 'sequencing':
        body_h = len(fd.get('items', [])) * 6 * mm + 2 * mm
    elif fmt == 'reason_evidence_table':
        body_h = 6 * mm + (fd.get('rows', 2) + 1) * 11 * mm + 2 * mm
    elif fmt == 'two_part_ab':
        body_h = len(fd.get('parts', [])) * (8.5 * 1.35 + 10 * mm) + 2 * mm
    elif fmt == 'draw_lines_matching':
        n_rows = max(len(fd.get('left_items', [])), len(fd.get('right_items', [])))
        body_h = n_rows * (_DLM_BOX_H + 1.5 * mm) + 2 * mm
    else:
        body_h = 13 * mm

    return ref_h + q_text_h + body_h + 3 * mm


# ===========================================================================
# Mark scheme — flowing renderers
# ===========================================================================

def _ms_award_text(marks, fmt='open_line', n_items=None):
    """
    Return format-specific award criteria lines matching SATs mark scheme conventions.
    n_items: number of statements/items in the format (for true_false_table).
    """
    n = n_items or 0

    if fmt == 'tick_one':
        return ["Award 1 mark for the correct answer."]

    if fmt == 'tick_two':
        return ["Award 1 mark for both correct answers."]

    if fmt == 'find_and_copy':
        return ["Award 1 mark for the correct word or phrase."]

    if fmt == 'sequencing':
        return ["Award 1 mark for the complete correct sequence."]

    if fmt == 'draw_lines_matching':
        return ["Award 1 mark for all lines joined to the correct boxes."]

    if fmt == 'true_false_table':
        if marks <= 1:
            return ([f"Award 1 mark for all {n} correct."] if n
                    else ["Award 1 mark for all correct."])
        else:
            return ([f"Award 1 mark for {n - 1} correct or 2 marks for all {n} correct."] if n
                    else ["Award 1 mark for all but one correct or 2 marks for all correct."])

    if fmt == 'two_part_ab':
        return ["Award 1 mark per part (up to 2 marks)."]

    # numbered_list: 1 mark = "write two things" (any correct answer earns the mark)
    if fmt == 'numbered_list' and marks == 1:
        return ["Award 1 mark for any correct answer (or both)."]

    # open_line, numbered_list (2+ marks), reason_evidence_table
    if marks == 1:
        return ["Award 1 mark for an acceptable answer."]
    if marks == 2:
        return ["Award 1 mark per acceptable point, up to a maximum of 2 marks."]
    if marks == 3:
        return [
            "Award 3 marks for two acceptable points, at least one with evidence.",
            "Award 2 marks for two acceptable points or one point with evidence.",
            "Award 1 mark for one acceptable point.",
        ]
    return [f"Award up to {marks} marks."]


def _draw_ms_award_header(c, q, y):
    """
    Draw award criteria line(s) in small italic grey above the answer.
    Called from render_ms_question() before every renderer.
    Returns new y.
    """
    fmt    = q.get('format', 'open_line')
    marks  = q.get('marks', 1)
    fd     = q.get('format_data', {})
    n_items = len(fd.get('statements', [])) if fmt == 'true_false_table' else None

    c.setFont("Helvetica-Oblique", 8)
    c.setFillColorRGB(0.3, 0.3, 0.3)
    for line in _ms_award_text(marks, fmt, n_items):
        for wrapped in wrap_text(c, line, "Helvetica-Oblique", 8, CW):
            c.drawString(MARGIN, y - 3.5 * mm, wrapped)
            y -= 8 * 1.35
    return y - 2 * mm


def render_ms_open_flowing(c, q, y):
    """Flowing mark scheme for open_line (2–3 marks), numbered_list, two_part_ab."""
    answer = _coerce_answer(q).strip()

    # "Accept:" header
    c.setFont("Helvetica-Bold", 8.5)
    c.setFillColorRGB(*GREEN)
    c.drawString(MARGIN, y - 3 * mm, "Accept:")
    y -= 8.5 * 1.4 + 1 * mm

    # Alternatives as bullets
    alternatives = [a.strip() for a in answer.replace('\n', '/').split('/') if a.strip()]
    if not alternatives:
        alternatives = [answer] if answer else ['\u2014']

    c.setFont("Helvetica", 8.5)
    c.setFillColorRGB(*GREEN)
    bullet = "\u2022  "
    bw = c.stringWidth(bullet, "Helvetica", 8.5)
    for alt in alternatives:
        c.drawString(MARGIN + 3 * mm, y - 3 * mm, bullet)
        alt_lines = wrap_text(c, alt, "Helvetica", 8.5, CW - 3 * mm - bw)
        for i, ln in enumerate(alt_lines):
            c.drawString(MARGIN + 3 * mm + bw, y - 3 * mm - i * (8.5 * 1.35), ln)
        y -= len(alt_lines) * (8.5 * 1.35) + 2 * mm

    # "Do not accept" block if present
    do_not = q.get('do_not_accept', '').strip()
    if do_not:
        y -= 2 * mm
        c.setFont("Helvetica-Bold", 8)
        c.setFillColorRGB(0.72, 0.18, 0.18)
        c.drawString(MARGIN, y - 3 * mm, "Do not accept:")
        y -= 8 * 1.4
        c.setFont("Helvetica-Oblique", 8)
        c.setFillColorRGB(0.72, 0.18, 0.18)
        for ln in wrap_text(c, do_not, "Helvetica-Oblique", 8, CW - 3 * mm):
            c.drawString(MARGIN + 3 * mm, y - 3 * mm, ln)
            y -= 8 * 1.35

    y -= 3 * mm
    c.setStrokeColorRGB(*GREY_LINE)
    c.setLineWidth(0.3)
    c.line(MARGIN, y, MARGIN + CW, y)
    return y - 4 * mm


def render_ms_reason_evidence_flowing(c, q, y):
    """Flowing mark scheme for reason_evidence_table questions."""
    answer = _coerce_answer(q).strip()

    # Parse "reason | evidence" rows
    rows = []
    for raw in answer.split('\n'):
        raw = raw.strip()
        if not raw:
            continue
        if '|' in raw:
            r, e = raw.split('|', 1)
            rows.append((r.strip(), e.strip()))
        else:
            rows.append((raw, ''))

    # "Acceptable points" header
    c.setFont("Helvetica-Bold", 8.5)
    c.setFillColorRGB(*BOX_BORDER)
    c.drawString(MARGIN, y - 3 * mm, "Acceptable points")
    y -= 8.5 * 1.4 + 1 * mm

    for i, (reason, evidence) in enumerate(rows, 1):
        # Numbered reason
        num_label = f"{i}.  "
        nw = c.stringWidth(num_label, "Helvetica-Bold", 8.5)
        c.setFont("Helvetica-Bold", 8.5)
        c.setFillColorRGB(*GREEN)
        c.drawString(MARGIN + 2 * mm, y - 3 * mm, num_label)
        reason_display = " / ".join(r.strip() for r in reason.split('/') if r.strip())
        r_lines = wrap_text(c, reason_display, "Helvetica-Bold", 8.5, CW - nw - 2 * mm)
        for j, ln in enumerate(r_lines):
            c.drawString(MARGIN + 2 * mm + nw, y - 3 * mm - j * (8.5 * 1.35), ln)
        y -= len(r_lines) * (8.5 * 1.35) + 1 * mm

        # Evidence bullets
        if evidence:
            ev_alts = [e.strip() for e in evidence.split('/') if e.strip()]
            bullet = "\u2022  "
            bw = c.stringWidth(bullet, "Helvetica", 8)
            c.setFont("Helvetica", 8)
            c.setFillColorRGB(*GREEN)
            for ev in ev_alts:
                c.drawString(MARGIN + 8 * mm, y - 2.5 * mm, bullet)
                ev_lines = wrap_text(c, ev, "Helvetica", 8, CW - 8 * mm - bw)
                for k, ln in enumerate(ev_lines):
                    c.drawString(MARGIN + 8 * mm + bw, y - 2.5 * mm - k * (8 * 1.35), ln)
                y -= len(ev_lines) * (8 * 1.35) + 1.5 * mm

        y -= 2 * mm

    c.setStrokeColorRGB(*GREY_LINE)
    c.setLineWidth(0.3)
    c.line(MARGIN, y, MARGIN + CW, y)
    return y - 4 * mm


def estimate_ms_height(c, q):
    """Estimate mark scheme height. More generous than pupil estimate."""
    fmt    = q.get('format', 'open_line')
    marks  = q.get('marks', 1)
    fd     = q.get('format_data', {})
    n_items = len(fd.get('statements', [])) if fmt == 'true_false_table' else None

    qtext = q.get('question', '')
    lw = c.stringWidth(f"{q['number']}. ", "Helvetica-Bold", 9)
    q_lines = wrap_text(c, qtext.split('\n\n')[-1], "Helvetica-Bold", 9, CW - lw)
    stem_h = len(q_lines) * 9 * 1.35 + 4 * mm

    # Award header height — present for every question type
    award_lines = len(_ms_award_text(marks, fmt, n_items))
    award_h = award_lines * 8 * 1.35 + 2 * mm

    # Simple formats — compact body + award header
    if fmt in ('tick_one', 'tick_two', 'true_false_table', 'find_and_copy',
               'sequencing', 'draw_lines_matching'):
        return stem_h + award_h + estimate_height(c, q) + 5 * mm

    if fmt == 'open_line' and marks == 1:
        return stem_h + award_h + 15 * mm

    answer = _coerce_answer(q).strip()

    if fmt == 'reason_evidence_table':
        rows = [r for r in answer.split('\n') if r.strip() and '|' in r]
        body_h = 10 * mm + len(rows) * 28 * mm + 8 * mm
    else:
        alternatives = [a.strip() for a in answer.replace('\n', '/').split('/') if a.strip()]
        n = max(len(alternatives), 1)
        body_h = 10 * mm + n * 8 * 1.35 * 1.8 + 8 * mm

    return stem_h + award_h + body_h + 10 * mm


def render_ms_question(c, q, y):
    """
    Render one question for the mark scheme PDF.
    Simple formats use compact answer renderers.
    Complex / open formats use flowing renderers.
    Returns new y, or None if question does not fit on this page.
    """
    if y - estimate_ms_height(c, q) < MIN_Y:
        return None

    y, _ = draw_preamble(c, q, y)
    marks_label(c, q.get('marks', 1), y + 9 * 1.35)

    # Award criteria — drawn for every question type
    y = _draw_ms_award_header(c, q, y)

    fmt   = q.get('format', 'open_line')
    marks = q.get('marks', 1)

    # Simple formats — compact renderers
    if fmt == 'tick_one':
        return render_tick_one_answer(c, q, y) - 3 * mm
    if fmt == 'tick_two':
        return render_tick_two_answer(c, q, y) - 3 * mm
    if fmt == 'true_false_table':
        return render_true_false_table_answer(c, q, y) - 3 * mm
    if fmt == 'find_and_copy':
        return render_find_and_copy_answer(c, q, y) - 3 * mm
    if fmt == 'sequencing':
        return render_sequencing_answer(c, q, y) - 3 * mm
    if fmt == 'draw_lines_matching':
        return render_draw_lines_matching_answer(c, q, y) - 3 * mm
    if fmt == 'open_line' and marks == 1:
        return render_open_line_answer(c, q, y) - 3 * mm

    # Complex / open formats — flowing, unconstrained
    if fmt == 'reason_evidence_table':
        return render_ms_reason_evidence_flowing(c, q, y)
    return render_ms_open_flowing(c, q, y)


# ===========================================================================
# Single question renderer — pupil / answer sheets (dispatches by format)
# ===========================================================================

def render_question(c, q, y, is_answer, is_supported):
    if y - estimate_height(c, q) < MIN_Y:
        return None

    y, _ = draw_preamble(c, q, y)
    marks_label(c, q.get('marks', 1), y + 9 * 1.35)

    fmt = q.get('format', 'open_line')

    if is_answer:
        dispatch = {
            'open_line':              lambda: render_open_line_answer(c, q, y),
            'find_and_copy':          lambda: render_find_and_copy_answer(c, q, y),
            'numbered_list':          lambda: render_numbered_list_answer(c, q, y),
            'tick_one':               lambda: render_tick_one_answer(c, q, y),
            'tick_two':               lambda: render_tick_two_answer(c, q, y),
            'true_false_table':       lambda: render_true_false_table_answer(c, q, y),
            'sequencing':             lambda: render_sequencing_answer(c, q, y),
            'reason_evidence_table':  lambda: render_reason_evidence_answer(c, q, y),
            'two_part_ab':            lambda: render_two_part_ab_answer(c, q, y),
            'draw_lines_matching':    lambda: render_draw_lines_matching_answer(c, q, y),
        }
    else:
        dispatch = {
            'open_line':              lambda: render_open_line_pupil(c, q, y, is_supported),
            'find_and_copy':          lambda: render_find_and_copy_pupil(c, q, y, is_supported),
            'numbered_list':          lambda: render_numbered_list_pupil(c, q, y, is_supported),
            'tick_one':               lambda: render_tick_one_pupil(c, q, y),
            'tick_two':               lambda: render_tick_two_pupil(c, q, y),
            'true_false_table':       lambda: render_true_false_table_pupil(c, q, y),
            'sequencing':             lambda: render_sequencing_pupil(c, q, y),
            'reason_evidence_table':  lambda: render_reason_evidence_pupil(c, q, y),
            'two_part_ab':            lambda: render_two_part_ab_pupil(c, q, y, is_supported),
            'draw_lines_matching':    lambda: render_draw_lines_matching_pupil(c, q, y),
        }

    renderer = dispatch.get(fmt)
    if renderer is None:
        return render_open_line_pupil(c, q, y, is_supported) - 3 * mm
    return renderer() - 3 * mm


# ===========================================================================
# Page builder
# ===========================================================================

def build_page(path, lesson, text, questions, is_answer, is_supported, icon_path):
    c = rl_canvas.Canvas(path, pagesize=A4)

    lesson_type = lesson['lesson_type'].capitalize()
    day = lesson['day']
    date = lesson['date']
    # Use learning_focus for the header; fall back to key_question for legacy data
    key_q = lesson.get('learning_focus') or lesson.get('key_question', '')
    i_cans = lesson.get('i_can_statements', [])

    y = draw_header(c, lesson_type, day, date, key_q, i_cans, icon_path)
    y = draw_text_box(c, text, y, font_size=10 if not is_supported else 9.5)

    for q in questions:
        result = render_question(c, q, y, is_answer=is_answer, is_supported=is_supported)
        if result is None:
            break
        y = result

    c.save()


def check_pages(path):
    return len(PdfReader(path).pages)


def merge_pdfs(paths, out_path):
    writer = PdfWriter()
    for p in paths:
        for page in PdfReader(p).pages:
            writer.add_page(page)
    with open(out_path, 'wb') as f:
        writer.write(f)


# ===========================================================================
# Text booklet page (SATs layout)
# ===========================================================================

def build_text_booklet_page(path, lesson, icon_path):
    c = rl_canvas.Canvas(path, pagesize=A4)

    lesson_type = lesson["lesson_type"].capitalize()
    day = lesson["day"]
    date_str = lesson["date"]
    key_q = lesson.get("learning_focus") or lesson.get("key_question", "")
    i_cans = lesson.get("i_can_statements", [])

    y = draw_header(c, lesson_type, day, date_str, key_q, i_cans, icon_path)

    text = lesson.get("text", "")
    if text:
        lines = wrap_text(c, text, "Helvetica", 11, CW - 6 * mm)
        lh = 11 * 1.5
        box_h = len(lines) * lh + 8 * mm

        c.setFillColorRGB(*BOX_BG)
        c.setStrokeColorRGB(*BOX_BORDER)
        c.setLineWidth(0.8)
        c.roundRect(MARGIN, y - box_h, CW, box_h, 2 * mm, fill=1, stroke=1)

        c.setFillColorRGB(*DARK)
        c.setFont("Helvetica", 11)
        ty = y - 4 * mm - 11 * 0.72
        for line in lines:
            c.drawString(MARGIN + 3 * mm, ty, line)
            ty -= lh

    c.save()


def build_pdfs(content: dict, icon_path: str, output_dir: str,
               layout: str = "integrated", year_group: str = "Y4") -> dict:
    import tempfile, shutil

    set_year_group(year_group)
    tmp = tempfile.mkdtemp()
    key_q = content.get('key_question', '')

    for lesson in content['lessons']:
        lesson['key_question'] = key_q

    std_pages, sup_pages, ans_pages, txt_pages = [], [], [], []

    for lesson in content['lessons']:
        lt = lesson['lesson_type']
        std_qs = lesson['questions']
        sup_qs = lesson['questions'][:5]
        text = lesson.get('text', '')

        if layout == 'sats':
            tp = os.path.join(tmp, f"{lt}_text.pdf")
            build_text_booklet_page(tp, lesson, icon_path)
            txt_pages.append(tp)

            p = os.path.join(tmp, f"{lt}_std.pdf")
            build_page(p, lesson, '', std_qs, is_answer=False,
                       is_supported=False, icon_path=icon_path)
            if check_pages(p) > 1:
                build_page(p, lesson, '', std_qs[:-1], is_answer=False,
                           is_supported=False, icon_path=icon_path)
            std_pages.append(p)

            p = os.path.join(tmp, f"{lt}_sup.pdf")
            build_page(p, lesson, '', sup_qs, is_answer=False,
                       is_supported=True, icon_path=icon_path)
            if check_pages(p) > 1:
                build_page(p, lesson, '', sup_qs[:-1], is_answer=False,
                           is_supported=True, icon_path=icon_path)
            sup_pages.append(p)

            p = os.path.join(tmp, f"{lt}_std_ans.pdf")
            build_page(p, lesson, '', std_qs, is_answer=True,
                       is_supported=False, icon_path=icon_path)
            p2 = os.path.join(tmp, f"{lt}_sup_ans.pdf")
            build_page(p2, lesson, '', sup_qs, is_answer=True,
                       is_supported=True, icon_path=icon_path)
            ans_pages.extend([p, p2])

        else:
            p = os.path.join(tmp, f"{lt}_std.pdf")
            build_page(p, lesson, text, std_qs, is_answer=False,
                       is_supported=False, icon_path=icon_path)
            if check_pages(p) > 1:
                build_page(p, lesson, text, std_qs[:-1], is_answer=False,
                           is_supported=False, icon_path=icon_path)
            std_pages.append(p)

            p = os.path.join(tmp, f"{lt}_sup.pdf")
            build_page(p, lesson, text, sup_qs, is_answer=False,
                       is_supported=True, icon_path=icon_path)
            if check_pages(p) > 1:
                build_page(p, lesson, text, sup_qs[:-1], is_answer=False,
                           is_supported=True, icon_path=icon_path)
            sup_pages.append(p)

            p = os.path.join(tmp, f"{lt}_std_ans.pdf")
            build_page(p, lesson, text, std_qs, is_answer=True,
                       is_supported=False, icon_path=icon_path)
            p2 = os.path.join(tmp, f"{lt}_sup_ans.pdf")
            build_page(p2, lesson, text, sup_qs, is_answer=True,
                       is_supported=True, icon_path=icon_path)
            ans_pages.extend([p, p2])

    os.makedirs(output_dir, exist_ok=True)
    out_std = os.path.join(output_dir, 'Standard_Pupil.pdf')
    out_sup = os.path.join(output_dir, 'Supported_Pupil.pdf')
    out_ans = os.path.join(output_dir, 'All_Answers.pdf')

    merge_pdfs(std_pages, out_std)
    merge_pdfs(sup_pages, out_sup)
    merge_pdfs(ans_pages, out_ans)

    result = {
        'standard': out_std,
        'supported': out_sup,
        'answers': out_ans,
    }

    if layout == 'sats' and txt_pages:
        out_txt = os.path.join(output_dir, 'Text_Booklet.pdf')
        merge_pdfs(txt_pages, out_txt)
        result['text_booklet'] = out_txt

    shutil.rmtree(tmp)
    return result


# ===========================================================================
# Structured text rendering (supports **heading** markers for non-fiction)
# ===========================================================================

def _parse_structured_text(text):
    segments = []
    for block in text.split('\n\n'):
        block = block.strip()
        if not block:
            continue
        if block.startswith('**') and block.endswith('**') and block.count('**') == 2:
            segments.append((True, block[2:-2].strip()))
        else:
            clean = block.replace('**', '')
            segments.append((False, clean))
    return segments or [(False, text)]


def draw_structured_text_box(c, text, y_top, font_size=10):
    segments = _parse_structured_text(text)

    lh_body = font_size * 1.42
    lh_head = font_size * 1.35
    gap = 2.5 * mm
    total_h = 5 * mm
    for is_heading, content in segments:
        font = "Helvetica-Bold" if is_heading else "Helvetica"
        lh = lh_head if is_heading else lh_body
        lines = wrap_text(c, content, font, font_size, CW - 6 * mm)
        total_h += len(lines) * lh + gap

    c.setFillColorRGB(*BOX_BG)
    c.setStrokeColorRGB(*BOX_BORDER)
    c.setLineWidth(0.8)
    c.roundRect(MARGIN, y_top - total_h, CW, total_h, 2 * mm, fill=1, stroke=1)

    ty = y_top - 3.5 * mm
    c.setFillColorRGB(*DARK)
    for is_heading, content in segments:
        font = "Helvetica-Bold" if is_heading else "Helvetica"
        lh = lh_head if is_heading else lh_body
        c.setFont(font, font_size)
        lines = wrap_text(c, content, font, font_size, CW - 6 * mm)
        for line in lines:
            ty -= lh
            c.drawString(MARGIN + 3 * mm, ty, line)
        ty -= gap

    return y_top - total_h - 3 * mm


# ===========================================================================
# Multi-page question rendering
# ===========================================================================

def _start_question_page(tmp_dir, prefix, page_num, header_kwargs, include_label):
    path = os.path.join(tmp_dir, f"{prefix}_p{page_num}.pdf")
    c = rl_canvas.Canvas(path, pagesize=A4)
    if include_label and page_num == 1:
        y = draw_header(c, **header_kwargs)
    else:
        y = H - MARGIN - 4 * mm
        if page_num > 1:
            c.setFont("Helvetica-Oblique", 7.5)
            c.setFillColorRGB(0.5, 0.5, 0.5)
            c.drawRightString(MARGIN + CW, y + 1 * mm, "continued…")
            y -= 5 * mm
    return c, y, path


def build_question_pages(
    tmp_dir, prefix, questions,
    is_answer, is_supported,
    icon_path, header_kwargs, include_label,
):
    pages = []
    page_num = 1
    c, y, path = _start_question_page(
        tmp_dir, prefix, page_num, header_kwargs, include_label
    )

    for q in questions:
        result = render_question(c, q, y, is_answer=is_answer, is_supported=is_supported)
        if result is None:
            c.save()
            pages.append(path)
            page_num += 1
            c, y, path = _start_question_page(
                tmp_dir, prefix, page_num, header_kwargs, include_label=False
            )
            result = render_question(c, q, y, is_answer=is_answer, is_supported=is_supported)
            if result is None:
                result = y - 25 * mm
        y = result

    c.save()
    pages.append(path)
    return pages


# ===========================================================================
# Reading Paper PDF builder
# ===========================================================================

def _build_reading_paper_text_pages(tmp_dir, text, title, font_size=11):
    """Build text-only page(s). Returns list of page paths."""
    lh_body  = font_size * 1.42
    lh_head  = font_size * 1.35
    seg_gap  = 2.5 * mm
    padding  = 3.5 * mm

    segments = _parse_structured_text(text)
    pages    = []
    seg_idx  = 0
    page_num = 0

    while seg_idx < len(segments):
        page_num += 1
        p = os.path.join(tmp_dir, f"rp_text_p{page_num}.pdf")
        c = rl_canvas.Canvas(p, pagesize=A4)

        if page_num == 1 and title:
            y = H - MARGIN - 2 * mm
            # Paper title — "Reading Paper: [topic]"
            paper_title = f"Reading Paper: {title}"
            c.setFont("Helvetica-Bold", 12)
            c.setFillColorRGB(*_ACCENT)
            for line in wrap_text(c, paper_title, "Helvetica-Bold", 12, CW):
                c.drawString(MARGIN, y, line)
                y -= 12 * 1.4
            y -= 3 * mm
            c.setStrokeColorRGB(*GREY_LINE)
            c.setLineWidth(0.5)
            c.line(MARGIN, y, MARGIN + CW, y)
            y -= 5 * mm
            box_top = y
        else:
            box_top = H - MARGIN - 2 * mm

        box_bottom = MARGIN

        c.setFillColorRGB(*BOX_BG)
        c.setStrokeColorRGB(*BOX_BORDER)
        c.setLineWidth(0.8)
        c.roundRect(MARGIN, box_bottom, CW, box_top - box_bottom,
                    2 * mm, fill=1, stroke=1)

        ty = box_top - padding
        c.setFillColorRGB(*DARK)

        while seg_idx < len(segments):
            is_heading, content = segments[seg_idx]
            font  = "Helvetica-Bold" if is_heading else "Helvetica"
            lh    = lh_head if is_heading else lh_body
            lines = wrap_text(c, content, font, font_size, CW - 6 * mm)
            seg_h = len(lines) * lh + seg_gap

            if ty - seg_h < box_bottom + padding:
                break

            # Keep heading with its following body — don't let it sit alone at page bottom
            if is_heading and seg_idx + 1 < len(segments):
                next_is_h, next_content = segments[seg_idx + 1]
                next_font = "Helvetica-Bold" if next_is_h else "Helvetica"
                next_lh   = lh_head if next_is_h else lh_body
                next_lines = wrap_text(c, next_content, next_font, font_size, CW - 6 * mm)
                # Require at least 2 lines of the next segment to fit on this page too
                min_follow_h = next_lh * min(2, len(next_lines)) + seg_gap
                if ty - seg_h - min_follow_h < box_bottom + padding:
                    break  # push heading onto next page with its content

            c.setFont(font, font_size)
            for line in lines:
                ty -= lh
                c.drawString(MARGIN + 3 * mm, ty, line)
            ty -= seg_gap
            seg_idx += 1

        c.save()
        pages.append(p)

    return pages


def build_reading_paper_pdfs(
    content: dict,
    icon_path: str,
    output_dir: str,
    include_label: bool = False,
    custom_label: str = "",
    year_group: str = "Y4",
) -> dict:
    """
    Build PDFs for Reading Paper Mode.
    content must have 'standard_text' and 'questions'.
    Returns dict: { 'text': path, 'questions': path, 'mark_scheme': path }
    """
    import tempfile, shutil

    set_year_group(year_group)
    tmp = tempfile.mkdtemp()
    topic    = content.get("key_question", "") or content.get("topic", "")
    text     = content.get("standard_text", "")
    questions = content.get("questions", [])

    # Shared title string used on both PDFs to link them
    paper_title = f"Reading Paper: {topic}" if topic else "Reading Paper"
    ms_title    = f"Mark Scheme: {topic}"   if topic else "Mark Scheme"

    # ── Text page(s) ──────────────────────────────────────────────────────
    txt_pages = _build_reading_paper_text_pages(tmp, text, topic)

    # ── Optional learning label renderer ──────────────────────────────────
    if include_label and custom_label.strip():
        def _draw_rp_label(c):
            y = H - MARGIN
            c.setFont("Helvetica-Bold", 8)
            c.setFillColorRGB(*DARK)
            c.drawString(MARGIN, y - 5 * mm, custom_label.strip())
            y -= 7 * mm
            c.setFont("Helvetica-Bold", 10)
            c.setFillColorRGB(*BOX_BORDER)
            c.drawString(MARGIN, y - 4 * mm, topic)
            kq_w = c.stringWidth(topic, "Helvetica-Bold", 10)
            c.setLineWidth(0.5)
            c.setStrokeColorRGB(*BOX_BORDER)
            c.line(MARGIN, y - 5 * mm, MARGIN + kq_w, y - 5 * mm)
            y -= 8 * mm
            c.setStrokeColorRGB(*GREY_LINE)
            c.setLineWidth(0.5)
            c.line(MARGIN, y, MARGIN + CW, y)
            y -= 3 * mm
            return y
    else:
        _draw_rp_label = None

    # ── Question pages ─────────────────────────────────────────────────────
    q_pages  = []
    page_num = 0

    def new_q_page(is_first):
        nonlocal page_num
        page_num += 1
        p = os.path.join(tmp, f"rp_q_p{page_num}.pdf")
        c = rl_canvas.Canvas(p, pagesize=A4)

        if is_first:
            y = H - MARGIN - 2 * mm
            # Paper title — left
            c.setFont("Helvetica-Bold", 11)
            c.setFillColorRGB(*_ACCENT)
            c.drawString(MARGIN, y, paper_title)
            # Name space — right, same line
            c.setFont("Helvetica", 9)
            c.setFillColorRGB(*DARK)
            c.drawRightString(MARGIN + CW, y, "Name:  ___________________________________")
            y -= 11 * 1.4 + 3 * mm
            c.setStrokeColorRGB(*GREY_LINE)
            c.setLineWidth(0.4)
            c.line(MARGIN, y, MARGIN + CW, y)
            y -= 4 * mm

            if _draw_rp_label:
                y = _draw_rp_label(c)
        else:
            y = H - MARGIN - 4 * mm
            c.setFont("Helvetica-Oblique", 7.5)
            c.setFillColorRGB(0.5, 0.5, 0.5)
            c.drawRightString(MARGIN + CW, y + 1 * mm, f"{paper_title} — continued…")
            y -= 5 * mm

        return c, y, p

    c, y, path = new_q_page(is_first=True)

    for q in questions:
        result = render_question(c, q, y, is_answer=False, is_supported=False)
        if result is None:
            c.save()
            q_pages.append(path)
            c, y, path = new_q_page(is_first=False)
            result = render_question(c, q, y, is_answer=False, is_supported=False)
            if result is None:
                result = y - 25 * mm
        y = result

    c.save()
    q_pages.append(path)

    # ── Mark scheme pages ──────────────────────────────────────────────────
    ms_pages   = []
    page_num_ms = 0

    def new_ms_page(is_first):
        nonlocal page_num_ms
        page_num_ms += 1
        p = os.path.join(tmp, f"rp_ms_p{page_num_ms}.pdf")
        c = rl_canvas.Canvas(p, pagesize=A4)
        y = H - MARGIN - 2 * mm

        if is_first:
            # Mark scheme title — links to paper_title
            c.setFont("Helvetica-Bold", 11)
            c.setFillColorRGB(*_ACCENT)
            c.drawString(MARGIN, y, ms_title)
            y -= 11 * 1.4
            # Subtitle: "for Reading Paper: [topic]"
            c.setFont("Helvetica", 8.5)
            c.setFillColorRGB(0.4, 0.4, 0.4)
            c.drawString(MARGIN, y, f"for {paper_title}")
            y -= 8.5 * 1.4 + 2 * mm
            c.setStrokeColorRGB(*GREY_LINE)
            c.setLineWidth(0.5)
            c.line(MARGIN, y, MARGIN + CW, y)
            y -= 5 * mm
        else:
            c.setFont("Helvetica-Oblique", 7.5)
            c.setFillColorRGB(0.5, 0.5, 0.5)
            c.drawRightString(MARGIN + CW, y + 1 * mm, f"{ms_title} — continued…")
            y -= 5 * mm

        return c, y, p

    c, y, path = new_ms_page(is_first=True)

    for q in questions:
        # Use render_ms_question — flowing for complex formats
        result = render_ms_question(c, q, y)
        if result is None:
            c.save()
            ms_pages.append(path)
            c, y, path = new_ms_page(is_first=False)
            result = render_ms_question(c, q, y)
            if result is None:
                result = y - 25 * mm
        y = result

    c.save()
    ms_pages.append(path)

    # ── Merge and write outputs ────────────────────────────────────────────
    os.makedirs(output_dir, exist_ok=True)
    out_text = os.path.join(output_dir, "ReadingPaper_Text.pdf")
    out_q    = os.path.join(output_dir, "ReadingPaper_Questions.pdf")
    out_ms   = os.path.join(output_dir, "ReadingPaper_MarkScheme.pdf")

    merge_pdfs(txt_pages, out_text)
    merge_pdfs(q_pages,   out_q)
    merge_pdfs(ms_pages,  out_ms)

    shutil.rmtree(tmp)

    return {
        "text":        out_text,
        "questions":   out_q,
        "mark_scheme": out_ms,
    }


# ===========================================================================
# KS1 PDF builders
# ===========================================================================

def _ks1_page_header(c, title, subtitle, is_first, is_ms=False):
    """Draw header for KS1 paper pages. Returns starting y."""
    y = H - MARGIN - 2 * mm
    if is_first:
        c.setFont("Helvetica-Bold", 11)
        c.setFillColorRGB(*_ACCENT)
        c.drawString(MARGIN, y, title)
        if subtitle:
            c.setFont("Helvetica", 9)
            c.setFillColorRGB(*DARK)
            c.drawRightString(MARGIN + CW, y, subtitle)
        y -= 11 * 1.4 + 1 * mm
        c.setStrokeColorRGB(*GREY_LINE)
        c.setLineWidth(0.4)
        c.line(MARGIN, y, MARGIN + CW, y)
        y -= 4 * mm
    else:
        c.setFont("Helvetica-Oblique", 7.5)
        c.setFillColorRGB(0.5, 0.5, 0.5)
        cont = f"{title} — continued…"
        c.drawRightString(MARGIN + CW, y + 1 * mm, cont)
        y -= 5 * mm
    return y


def _ks1_passage_title_band(c, passage_title, text_type_label, y):
    """Draw a coloured band marking the start of a new passage."""
    band_h = 8 * mm
    c.setFillColorRGB(*_ACCENT)
    c.rect(MARGIN, y - band_h, CW, band_h, fill=1, stroke=0)
    c.setFillColorRGB(1, 1, 1)
    c.setFont("Helvetica-Bold", 10)
    c.drawString(MARGIN + 3 * mm, y - band_h + 2.5 * mm, passage_title)
    if text_type_label:
        c.setFont("Helvetica", 8)
        c.drawRightString(MARGIN + CW - 3 * mm, y - band_h + 2.5 * mm, text_type_label)
    return y - band_h - 3 * mm


def _ks1_questions_for_passage(c, questions, y, passage_title,
                                is_answer, is_supported, tmp, pages):
    """
    Render a flat list of questions for a passage, handling page breaks.
    Returns final y.
    """
    page_num = [len(pages) + 1]

    def new_page():
        c_new_path = os.path.join(tmp, f"ks1_q_p{page_num[0]}.pdf")
        page_num[0] += 1
        c_new = rl_canvas.Canvas(c_new_path, pagesize=A4)
        y_new = _ks1_page_header(c_new, passage_title, "", is_first=False)
        return c_new, y_new, c_new_path

    for q in questions:
        result = render_question(c, q, y, is_answer=is_answer, is_supported=is_supported)
        if result is None:
            c.save()
            pages.append(None)  # placeholder — caller must track path
            c, y, _ = new_page()
            result = render_question(c, q, y, is_answer=is_answer, is_supported=is_supported)
            if result is None:
                result = y - 25 * mm
        y = result

    return c, y


def build_ks1_paper2_pdfs(
    content: dict,
    output_dir: str,
    year_group: str = "Y2",
) -> dict:
    """
    Build KS1 Paper 2 PDFs (separate text booklet + answer booklet).
    content: output of generate_ks1_paper() with paper_type="separate"
    Returns: { "text": path, "questions": path, "supported": path, "mark_scheme": path }
    """
    import tempfile, shutil as _shutil

    set_year_group(year_group)
    tmp = tempfile.mkdtemp()

    p1 = content["passage1"]
    p2 = content["passage2"]

    p1_title = p1.get("key_question", "") or "Passage 1"
    p2_title = p2.get("key_question", "") or "Passage 2"

    paper_label = f"KS1 Reading Paper ({year_group})"
    ms_label    = f"Mark Scheme — KS1 Reading Paper ({year_group})"

    # ── Text booklet ──────────────────────────────────────────────────────────
    txt_pages = []
    for idx, (passage, ptitle) in enumerate([(p1, p1_title), (p2, p2_title)], 1):
        txt_pages += _build_reading_paper_text_pages(
            tmp, passage.get("standard_text", ""),
            ptitle, font_size=11,
        )

    # ── Answer booklet (standard) ─────────────────────────────────────────────
    def _build_answer_booklet(is_supported):
        pages_out  = []
        page_count = [0]

        def new_page(is_first):
            page_count[0] += 1
            p = os.path.join(tmp, f"ks1_{'sup' if is_supported else 'std'}_{page_count[0]}.pdf")
            c = rl_canvas.Canvas(p, pagesize=A4)
            sub = "Name: ___________________________" if is_first else ""
            y = _ks1_page_header(c, paper_label, sub, is_first=is_first)
            return c, y, p

        c, y, path = new_page(is_first=True)

        for passage, ptitle, all_qs in [
            (p1, p1_title, p1.get("questions", [])),
            (p2, p2_title, p2.get("questions", [])),
        ]:
            # Passage title band
            if y < MARGIN + 40 * mm:
                c.save()
                pages_out.append(path)
                c, y, path = new_page(is_first=False)

            y = _ks1_passage_title_band(c, ptitle, "", y)

            for q in all_qs:
                result = render_question(c, q, y,
                                         is_answer=False,
                                         is_supported=is_supported)
                if result is None:
                    c.save()
                    pages_out.append(path)
                    c, y, path = new_page(is_first=False)
                    result = render_question(c, q, y,
                                             is_answer=False,
                                             is_supported=is_supported)
                    if result is None:
                        result = y - 25 * mm
                y = result

        c.save()
        pages_out.append(path)
        return pages_out

    std_pages = _build_answer_booklet(is_supported=False)
    sup_pages = _build_answer_booklet(is_supported=True)

    # ── Mark scheme ───────────────────────────────────────────────────────────
    ms_pages   = []
    ms_count   = [0]

    def new_ms_page(is_first):
        ms_count[0] += 1
        p = os.path.join(tmp, f"ks1_ms_{ms_count[0]}.pdf")
        c = rl_canvas.Canvas(p, pagesize=A4)
        y = _ks1_page_header(c, ms_label, f"for {paper_label}", is_first=is_first)
        return c, y, p

    c, y, path = new_ms_page(is_first=True)

    for passage, ptitle, all_qs in [
        (p1, p1_title, p1.get("questions", [])),
        (p2, p2_title, p2.get("questions", [])),
    ]:
        if y < MARGIN + 40 * mm:
            c.save()
            ms_pages.append(path)
            c, y, path = new_ms_page(is_first=False)
        y = _ks1_passage_title_band(c, ptitle, "", y)

        for q in all_qs:
            result = render_ms_question(c, q, y)
            if result is None:
                c.save()
                ms_pages.append(path)
                c, y, path = new_ms_page(is_first=False)
                result = render_ms_question(c, q, y)
                if result is None:
                    result = y - 25 * mm
            y = result

    c.save()
    ms_pages.append(path)

    # ── Merge ─────────────────────────────────────────────────────────────────
    os.makedirs(output_dir, exist_ok=True)
    out_text = os.path.join(output_dir, "KS1_Text.pdf")
    out_std  = os.path.join(output_dir, "KS1_Questions.pdf")
    out_sup  = os.path.join(output_dir, "KS1_Questions_Supported.pdf")
    out_ms   = os.path.join(output_dir, "KS1_MarkScheme.pdf")

    merge_pdfs(txt_pages, out_text)
    merge_pdfs(std_pages, out_std)
    merge_pdfs(sup_pages, out_sup)
    merge_pdfs(ms_pages,  out_ms)

    _shutil.rmtree(tmp)
    return {"text": out_text, "questions": out_std,
            "supported": out_sup, "mark_scheme": out_ms}


def build_ks1_paper1_pdfs(
    content: dict,
    output_dir: str,
    year_group: str = "Y2",
) -> dict:
    """
    Build KS1 Paper 1 PDFs (combined text + questions booklet).
    content: output of generate_ks1_paper() with paper_type="combined"
    Returns: { "combined": path, "mark_scheme": path }
    """
    import tempfile, shutil as _shutil

    set_year_group(year_group)
    tmp = tempfile.mkdtemp()

    p1 = content["passage1"]
    p2 = content["passage2"]

    paper_label = f"KS1 Reading Paper ({year_group})"
    ms_label    = f"Mark Scheme — KS1 Reading Paper ({year_group})"

    def _build_combined(is_ms):
        pages_out = []
        pg_count  = [0]

        def new_page(is_first):
            pg_count[0] += 1
            p = os.path.join(tmp, f"ks1_p1_{'ms' if is_ms else 'std'}_{pg_count[0]}.pdf")
            c = rl_canvas.Canvas(p, pagesize=A4)
            title = ms_label if is_ms else paper_label
            sub = ("" if is_ms else
                   "Name: ___________________________" if is_first else "")
            y = _ks1_page_header(c, title, sub, is_first=is_first)
            return c, y, p

        c, y, path = new_page(is_first=True)

        for passage_idx, passage in enumerate([p1, p2]):
            ptitle = passage.get("title", "")
            ptype  = passage.get("text_type", "")
            type_label = {"fiction": "Fiction", "non_fiction": "Non-fiction",
                          "poetry": "Poetry"}.get(ptype, "")

            # Force a new page before passage 2 — always
            if passage_idx == 1:
                c.save()
                pages_out.append(path)
                c, y, path = new_page(is_first=False)
            elif y < MARGIN + 50 * mm:
                c.save()
                pages_out.append(path)
                c, y, path = new_page(is_first=False)

            y = _ks1_passage_title_band(c, ptitle, type_label, y)

            for section in passage.get("sections", []):
                chunk = section.get("text_chunk", "")
                qs    = section.get("questions", [])

                # Rough check: if remaining space is under 80mm, start a new page
                if y < MARGIN + 80 * mm:
                    c.save()
                    pages_out.append(path)
                    c, y, path = new_page(is_first=False)

                # Draw text chunk in a box
                if chunk.strip():
                    y = draw_text_box(c, chunk, y_top=y, font_size=10)
                    y -= 2 * mm

                # Draw questions
                for q in qs:
                    # Paper 1 combined: suppress text_reference display
                    q_copy = dict(q)
                    q_copy['text_reference'] = ''

                    if is_ms:
                        result = render_ms_question(c, q_copy, y)
                    else:
                        result = render_question(c, q_copy, y,
                                                  is_answer=False, is_supported=False)
                    if result is None:
                        c.save()
                        pages_out.append(path)
                        c, y, path = new_page(is_first=False)
                        if is_ms:
                            result = render_ms_question(c, q_copy, y)
                        else:
                            result = render_question(c, q_copy, y,
                                                      is_answer=False, is_supported=False)
                        if result is None:
                            result = y - 25 * mm
                    y = result

                y -= 3 * mm  # small gap between sections

        c.save()
        pages_out.append(path)
        return pages_out

    combined_pages = _build_combined(is_ms=False)
    ms_pages       = _build_combined(is_ms=True)

    os.makedirs(output_dir, exist_ok=True)
    out_combined = os.path.join(output_dir, "KS1_Combined.pdf")
    out_ms       = os.path.join(output_dir, "KS1_MarkScheme.pdf")

    merge_pdfs(combined_pages, out_combined)
    merge_pdfs(ms_pages,       out_ms)

    _shutil.rmtree(tmp)
    return {"combined": out_combined, "mark_scheme": out_ms}


# ===========================================================================
# Decodable Reader output builders
# ===========================================================================

def build_decodable_pdf(result: dict, output_dir: str) -> str:
    """Build a clean A4 PDF for a decodable reading text."""
    import tempfile as _tmpmod, shutil as _shutil

    text        = result.get("text", "")
    phase       = result.get("phase", "")
    stage_label = result.get("stage_label", "")
    validation  = result.get("validation", {})
    flag_hard   = validation.get("flag_hard", [])
    flag_vce    = validation.get("flag_vce", [])

    set_year_group("Y1")

    tmp = _tmpmod.mkdtemp()
    p   = os.path.join(tmp, "dec_text.pdf")
    c   = rl_canvas.Canvas(p, pagesize=A4)

    band_h = 14 * mm
    c.setFillColorRGB(*_ACCENT)
    c.rect(0, H - band_h, W, band_h, fill=1, stroke=0)
    c.setFillColorRGB(1, 1, 1)
    c.setFont("Helvetica-Bold", 11)
    c.drawString(MARGIN, H - band_h + 4 * mm, phase)
    c.setFont("Helvetica", 9)
    c.drawRightString(W - MARGIN, H - band_h + 4 * mm, stage_label)

    y         = H - band_h - 10 * mm
    font_size = 18
    leading   = font_size * 1.5

    c.setFont("Helvetica", font_size)
    c.setFillColorRGB(*DARK)

    paragraphs = [p2.strip() for p2 in text.split('\n') if p2.strip()]
    for para in paragraphs:
        lines = wrap_text(c, para, "Helvetica", font_size, CW)
        for line in lines:
            if y < MARGIN + leading:
                c.showPage()
                c.setFillColorRGB(*DARK)
                c.setFont("Helvetica", font_size)
                y = H - MARGIN - leading
            c.drawString(MARGIN, y, line)
            y -= leading
        y -= leading * 0.4

    if flag_hard or flag_vce:
        y -= 4 * mm
        c.setStrokeColorRGB(*GREY_LINE)
        c.setLineWidth(0.4)
        c.line(MARGIN, y, MARGIN + CW, y)
        y -= 5 * mm
        c.setFont("Helvetica-Bold", 8)
        if flag_hard:
            c.setFillColorRGB(0.65, 0.1, 0.1)
            msg = f"Check these words (may be outside {phase} GPCs): {', '.join(flag_hard)}"
            for ln in wrap_text(c, msg, "Helvetica-Bold", 8, CW):
                c.drawString(MARGIN, y, ln)
                y -= 8 * 1.35
        if flag_vce:
            c.setFillColorRGB(0.55, 0.35, 0.0)
            msg = f"Possible split digraphs - check if intended at {phase}: {', '.join(flag_vce)}"
            for ln in wrap_text(c, msg, "Helvetica-Bold", 8, CW):
                c.drawString(MARGIN, y, ln)
                y -= 8 * 1.35

    c.save()
    os.makedirs(output_dir, exist_ok=True)
    out = os.path.join(output_dir, "Decodable_Text.pdf")
    import shutil
    shutil.copy(p, out)
    _shutil.rmtree(tmp)
    return out


def build_decodable_pptx(result: dict, output_dir: str) -> str:
    """Build a single-slide PPTX for displaying a decodable reading text."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    text        = result.get("text", "")
    phase       = result.get("phase", "")
    stage_label = result.get("stage_label", "")
    validation  = result.get("validation", {})
    flag_hard   = validation.get("flag_hard", [])
    flag_vce    = validation.get("flag_vce",  [])

    set_year_group("Y1")
    r, g, b = [int(v * 255) for v in _ACCENT]
    accent_rgb = RGBColor(r, g, b)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    band_h = Inches(0.45)
    band   = slide.shapes.add_shape(1, 0, 0, prs.slide_width, band_h)
    band.fill.solid()
    band.fill.fore_color.rgb = accent_rgb
    band.line.fill.background()
    tf = band.text_frame
    tf.text = f"{phase}  -  {stage_label}"
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    margin = Inches(0.55)
    top    = band_h + Inches(0.3)
    body_h = prs.slide_height - top - Inches(0.5)
    txBox  = slide.shapes.add_textbox(margin, top, prs.slide_width - 2 * margin, body_h)
    tf2    = txBox.text_frame
    tf2.word_wrap = True

    paragraphs = [p2.strip() for p2 in text.split('\n') if p2.strip()]
    for i, para_text in enumerate(paragraphs):
        p_obj = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p_obj.space_after  = Pt(8)
        run = p_obj.add_run()
        run.text = para_text
        run.font.size  = Pt(28)
        run.font.color.rgb = RGBColor(34, 34, 34)

    if flag_hard or flag_vce:
        note_top = prs.slide_height - Inches(0.65)
        note_box = slide.shapes.add_textbox(
            margin, note_top, prs.slide_width - 2 * margin, Inches(0.55)
        )
        nf   = note_box.text_frame
        nf.word_wrap = True
        parts = []
        if flag_hard:
            parts.append(f"Check: {', '.join(flag_hard)}")
        if flag_vce:
            parts.append(f"Split digraphs?: {', '.join(flag_vce)}")
        nr = nf.paragraphs[0].add_run()
        nr.text = "  |  ".join(parts)
        nr.font.size   = Pt(10)
        nr.font.italic = True
        nr.font.color.rgb = RGBColor(160, 60, 60)

    os.makedirs(output_dir, exist_ok=True)
    out = os.path.join(output_dir, "Decodable_Text.pptx")
    prs.save(out)
    return out
